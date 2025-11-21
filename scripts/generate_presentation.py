#!/usr/bin/env python3
"""Genera presentación ejecutiva (HU-PRESENT-01) en Markdown y opcional PDF/PPTX.
Estructura: Contexto, Dataset, Metodología, Resultados (clasificación, regresión, clustering, anomalías), Interpretabilidad, Recomendaciones.
Uso:
  python3 scripts/generate_presentation.py [--export-pdf] [--export-pptx]
Requisitos opcionales:
  - pypandoc o reportlab para PDF
  - python-pptx para PPTX
"""
import argparse, os, json, datetime
from pathlib import Path
import pandas as pd

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / 'docs'
DOCS.mkdir(exist_ok=True)
OUT_MD = DOCS / 'presentacion.md'

# Artefactos
paths = {
    'metrics_clf': ROOT / 'reports' / 'metrics_classification.txt',
    'metrics_reg': ROOT / 'reports' / 'metrics_regression.txt',
    'feature_imp_clf': ROOT / 'reports' / 'feature_importance_classification.csv',
    'feature_imp_reg': ROOT / 'reports' / 'feature_importance_regression.csv',
    'hpo_summary': ROOT / 'reports' / 'hpo_summary.md',
    'clustering_results': ROOT / 'reports' / 'clustering_results.csv',
    'clustering_fig_km': ROOT / 'reports' / 'clustering_kmeans_silhouette.png',
    'clustering_fig_db': ROOT / 'reports' / 'clustering_dbscan_eps_grid.png',
    'anomaly_results': ROOT / 'reports' / 'anomaly_results.csv',
    'anomaly_fig_frac': ROOT / 'reports' / 'anomaly_fraction_by_algo.png',
    'leakage_report': ROOT / 'reports' / 'leakage_report.json',
    'correlation_matrix': ROOT / 'data' / 'processed' / 'correlation_matrix.csv',
}

now_iso = datetime.datetime.utcnow().isoformat() + 'Z'


def safe_read_text(p: Path) -> str:
    return p.read_text(encoding='utf-8').strip() if p.exists() else 'pendiente'

def load_csv_head(p: Path, n=5):
    if not p.exists():
        return 'pendiente'
    try:
        df = pd.read_csv(p)
        return df.head(n).to_csv(index=False, sep='|')
    except Exception:
        return 'pendiente'

# Resultados Clustering
clustering_section = 'pendiente'
if paths['clustering_results'].exists():
    try:
        cdf = pd.read_csv(paths['clustering_results'])
        top = cdf.sort_values('silhouette', ascending=False).head(3)
        clustering_section = '\n'.join([
            'Top clustering (silhouette):',
            top.to_csv(index=False, sep='|')
        ])
    except Exception:
        pass

# Resultados Anomalías
anomaly_section = 'pendiente'
if paths['anomaly_results'].exists():
    try:
        adf = pd.read_csv(paths['anomaly_results'])
        sel = adf[['algo','anomaly_fraction','contamination_cfg','runtime_ms']].copy().head(5)
        anomaly_section = '\n'.join([
            'Anomalías (fracción vs target):',
            sel.to_csv(index=False, sep='|')
        ])
    except Exception:
        pass

# Interpretabilidad
interpret_section = []
if paths['feature_imp_clf'].exists():
    try:
        df_fi = pd.read_csv(paths['feature_imp_clf']).sort_values('importance', ascending=False).head(5)
        interpret_section.append('Top importancias clasificación:')
        interpret_section.append(df_fi.to_csv(index=False, sep='|'))
    except Exception:
        interpret_section.append('Importancias clasificación: pendiente')
else:
    interpret_section.append('Importancias clasificación: pendiente')
if paths['feature_imp_reg'].exists():
    try:
        df_fi_r = pd.read_csv(paths['feature_imp_reg']).sort_values('importance', ascending=False).head(5)
        interpret_section.append('Top importancias regresión:')
        interpret_section.append(df_fi_r.to_csv(index=False, sep='|'))
    except Exception:
        interpret_section.append('Importancias regresión: pendiente')
else:
    interpret_section.append('Importancias regresión: pendiente')
interpret_md = '\n'.join(interpret_section)

# Métricas básicas
metrics_clf = safe_read_text(paths['metrics_clf'])
metrics_reg = safe_read_text(paths['metrics_reg'])

# Leakage
leakage_line = 'pendiente'
if paths['leakage_report'].exists():
    try:
        leak = json.loads(paths['leakage_report'].read_text(encoding='utf-8'))
        leakage_line = f"Leakage flag={leak.get('flag')} r2={leak.get('r2')} features={leak.get('tested_features')}"
    except Exception:
        pass

# Correlation size
corr_line = 'pendiente'
if paths['correlation_matrix'].exists():
    try:
        corr_df = pd.read_csv(paths['correlation_matrix'], index_col=0)
        corr_line = f"Matriz de correlación: {corr_df.shape[0]} variables"
    except Exception:
        pass

hpo_line = 'pendiente'
if paths['hpo_summary'].exists():
    hpo_line = 'HPO ejecutado (ver hpo_summary.md para top configuraciones).'

contexto = f"""## Contexto
Problema: Comprender y modelar modalidad de programas y edad promedio, explorando también estructuras no supervisadas (clustering) y anomalías para calidad de datos.
Fecha de generación: {now_iso}
Leakage: {leakage_line}
"""

dataset = f"""## Dataset
{corr_line}
Fuente: data/raw/*.csv (ver config/config.py para parámetros). Tamaño aproximado post-proceso: ver X_train_engineered.csv.
"""

metodologia = f"""## Metodología
Pipeline: EDA -> Preprocesamiento -> Feature Engineering -> Entrenamiento -> Evaluación -> Interpretabilidad -> Clustering -> Anomalías.
Validación temporal aplicada (split por año). HPO: {hpo_line}
"""

resultados = f"""## Resultados
### Clasificación
Métrica: {metrics_clf}
Figura/Importancias: ver reports/feature_importance_classification.csv (top 5 abajo si disponible).

### Regresión
Métrica: {metrics_reg}
Figura/Importancias: ver reports/feature_importance_regression.csv.

### Clustering
{clustering_section}
Figura: {paths['clustering_fig_db'].name if paths['clustering_fig_db'].exists() else 'pendiente'}

### Anomalías
{anomaly_section}
Figura: {paths['anomaly_fig_frac'].name if paths['anomaly_fig_frac'].exists() else 'pendiente'}
"""

interpretabilidad = f"""## Interpretabilidad
{interpret_md}
"""

recomendaciones = """## Recomendaciones
- Consolidar almacenamiento de modelos entrenados y versionar.
- Incorporar monitoreo de drift y recalibración anual.
- Extender modelos a boosting y SHAP para interpretabilidad avanzada.
- Integrar panel UI para exploración de anomalías y clusters.
"""

content = "\n".join([contexto, dataset, metodologia, resultados, interpretabilidad, recomendaciones])
OUT_MD.write_text(content, encoding='utf-8')
print(f"Presentación escrita en {OUT_MD}")


def export_pdf(md_path: Path):
    try:
        import pypandoc
        out_pdf = md_path.with_suffix('.pdf')
        pypandoc.convert_file(str(md_path), 'pdf', outputfile=str(out_pdf))
        print(f"PDF exportado: {out_pdf}")
    except Exception as e:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas
            out_pdf = md_path.with_suffix('.pdf')
            c = canvas.Canvas(str(out_pdf), pagesize=A4)
            text = md_path.read_text(encoding='utf-8')
            y = 800
            import re, os
            from reportlab.lib.utils import ImageReader
            for line in text.splitlines():
                stripped=line.strip()
                img_match=re.match(r'!\[[^\]]*\]\(([^)]+)\)', stripped)
                if img_match:
                    raw_path=img_match.group(1)
                    img_path=Path(raw_path)
                    if not img_path.is_absolute():
                        img_path=(md_path.parent/Path(raw_path)).resolve()
                    if img_path.exists():
                        try:
                            im=ImageReader(str(img_path))
                            iw,ih=im.getSize()
                            max_w=500; scale=max_w/iw; dw=iw*scale; dh=ih*scale
                            if y - dh < 40:
                                c.showPage(); y=800
                            c.drawImage(im,40,y-dh,dw,dh,preserveAspectRatio=True,anchor='nw')
                            y -= (dh+12)
                            continue
                        except Exception:
                            pass
                c.drawString(40, y, line[:120])
                y -= 14
                if y < 40:
                    c.showPage(); y = 800
            c.save()
            print(f"PDF exportado (fallback reportlab): {out_pdf}")
        except Exception as e2:
            print(f"[WARN] No se pudo exportar PDF: {e2}")


def export_pptx(md_path: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        sections_titles = ["Contexto","Dataset","Metodología","Resultados - Clasificación","Resultados - Regresión","Resultados - Clustering","Resultados - Anomalías","Interpretabilidad","Recomendaciones"]
        md_text = md_path.read_text(encoding='utf-8')
        blocks = md_text.split('\n## ')
        # Simple mapping by title order
        for title in sections_titles:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            body = slide.shapes.placeholders[1].text_frame
            # Extract content lines matching title
            for blk in blocks:
                if blk.startswith(title):
                    lines = blk.splitlines()[1:][:12]
                    for i,l in enumerate(lines):
                        p = body.add_paragraph() if i>0 else body.paragraphs[0]
                        p.text = l[:180]
                        p.font.size = Pt(14)
                    break
        out_pptx = md_path.with_suffix('.pptx')
        prs.save(str(out_pptx))
        print(f"PPTX exportado: {out_pptx}")
    except Exception as e:
        print(f"[WARN] No se pudo exportar PPTX: {e}")

if __name__ == '__main__':
    ap = argparse.ArgumentParser(description='Generar presentación ejecutiva Markdown/PDF/PPTX')
    ap.add_argument('--export-pdf', action='store_true')
    ap.add_argument('--export-pptx', action='store_true')
    args = ap.parse_args()
    if args.export_pdf:
        export_pdf(OUT_MD)
    if args.export_pptx:
        export_pptx(OUT_MD)
